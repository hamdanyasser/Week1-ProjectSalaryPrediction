"""PayScope — 7-slide tight, vibrant, animated, chart-rich deck."""
from __future__ import annotations

from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# ── palette (matches dashboard) ─────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0E, 0x1A)
SURFACE  = RGBColor(0x12, 0x17, 0x2B)
SURF_ALT = RGBColor(0x19, 0x1F, 0x38)
TEAL     = RGBColor(0x00, 0xE5, 0xC3)
CORAL    = RGBColor(0xFF, 0x6B, 0x6B)
AMBER    = RGBColor(0xFF, 0xB5, 0x47)
VIOLET   = RGBColor(0xA7, 0x8B, 0xFA)
SKY      = RGBColor(0x38, 0xBD, 0xF8)
LIME     = RGBColor(0x84, 0xCC, 0x16)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
WHITE    = RGBColor(0xED, 0xF2, 0xF7)
MUTED    = RGBColor(0x8B, 0x95, 0xA8)
DARK     = RGBColor(0x0A, 0x0E, 0x1A)

W, H = Inches(13.333), Inches(7.5)

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


# ── low-level helpers ────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    if s.adjustments and len(s.adjustments) > 0:
        s.adjustments[0] = radius
    return s


def hard_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def circle(slide, l, t, size, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(s, alpha)
    return s


def _set_shape_alpha(shape, alpha_pct):
    spPr = shape._element.spPr
    solidFill = spPr.find(".//a:solidFill", NSMAP)
    if solidFill is not None:
        srgb = solidFill.find("a:srgbClr", NSMAP)
        if srgb is not None:
            for e in srgb.findall("a:alpha", NSMAP):
                srgb.remove(e)
            alpha_el = etree.SubElement(srgb, f"{{{NSMAP['a']}}}alpha")
            alpha_el.set("val", str(int((100 - alpha_pct) * 1000)))


def txt(slide, l, t, w, h, text, sz, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def pill(slide, l, t, text, bg_color, txt_color=DARK, w=Inches(2.0), h=Inches(0.42)):
    s = rect(slide, l, t, w, h, bg_color, radius=0.5)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Pt(8)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = txt_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    return s


def glow(slide, cx, cy, r, color, alpha=90):
    return circle(slide, cx - r, cy - r, r * 2, color, alpha)


def dot_grid(slide, x, y, cols, rows, gap, color, alpha=88, size=Inches(0.06)):
    for r in range(rows):
        for c in range(cols):
            circle(slide, x + c * gap, y + r * gap, size, color, alpha)


def accent_lines(slide, x, y, colors, w=Inches(0.5), h=Inches(0.08), gap=Inches(0.1)):
    for i, col in enumerate(colors):
        hard_rect(slide, x + i * (w + gap), y, w, h, col)


# ── animation ────────────────────────────────────────────────────────────

def anim(slide, shape, delay_ms=0, dur_ms=500, effect="fade"):
    """Entrance animation (fade / fly / zoom)."""
    slide_element = slide._element
    timing = slide_element.find(f"{{{NSMAP['p']}}}timing")
    if timing is None:
        timing = etree.SubElement(slide_element, f"{{{NSMAP['p']}}}timing")

    tnLst = timing.find(f"{{{NSMAP['p']}}}tnLst")
    if tnLst is None:
        tnLst = etree.SubElement(timing, f"{{{NSMAP['p']}}}tnLst")

    par = tnLst.find(f"{{{NSMAP['p']}}}par")
    if par is None:
        par = etree.SubElement(tnLst, f"{{{NSMAP['p']}}}par")
        cTn_root = etree.SubElement(par, f"{{{NSMAP['p']}}}cTn", attrib={
            "id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot"})
        childTnLst = etree.SubElement(cTn_root, f"{{{NSMAP['p']}}}childTnLst")
        seq = etree.SubElement(childTnLst, f"{{{NSMAP['p']}}}seq", attrib={
            "concurrent": "1", "nextAc": "seek"})
        seq_cTn = etree.SubElement(seq, f"{{{NSMAP['p']}}}cTn", attrib={
            "id": "2", "dur": "indefinite", "nodeType": "mainSeq"})
        seq_childTnLst = etree.SubElement(seq_cTn, f"{{{NSMAP['p']}}}childTnLst")
    else:
        cTn_root = par.find(f"{{{NSMAP['p']}}}cTn")
        childTnLst = cTn_root.find(f"{{{NSMAP['p']}}}childTnLst")
        seq = childTnLst.find(f"{{{NSMAP['p']}}}seq")
        seq_cTn = seq.find(f"{{{NSMAP['p']}}}cTn")
        seq_childTnLst = seq_cTn.find(f"{{{NSMAP['p']}}}childTnLst")

    existing = seq_childTnLst.findall(f"{{{NSMAP['p']}}}par")
    next_id = 3 + len(existing) * 10

    outer_par = etree.SubElement(seq_childTnLst, f"{{{NSMAP['p']}}}par")
    outer_cTn = etree.SubElement(outer_par, f"{{{NSMAP['p']}}}cTn", attrib={
        "id": str(next_id), "fill": "hold"})
    stCondLst = etree.SubElement(outer_cTn, f"{{{NSMAP['p']}}}stCondLst")
    cond = etree.SubElement(stCondLst, f"{{{NSMAP['p']}}}cond", attrib={"delay": "0"})
    if len(existing) == 0:
        cond.set("evt", "onNext")
        tgtEl = etree.SubElement(cond, f"{{{NSMAP['p']}}}tgtEl")
        etree.SubElement(tgtEl, f"{{{NSMAP['p']}}}sldTgt")

    outer_childTnLst = etree.SubElement(outer_cTn, f"{{{NSMAP['p']}}}childTnLst")

    inner_par = etree.SubElement(outer_childTnLst, f"{{{NSMAP['p']}}}par")
    preset_id = {"fade": "10", "fly": "2", "zoom": "23", "wipe": "12"}.get(effect, "10")
    inner_cTn = etree.SubElement(inner_par, f"{{{NSMAP['p']}}}cTn", attrib={
        "id": str(next_id + 1), "presetID": preset_id,
        "presetClass": "entr", "presetSubtype": "0",
        "fill": "hold", "grpId": "0", "nodeType": "withEffect"})
    inner_stCondLst = etree.SubElement(inner_cTn, f"{{{NSMAP['p']}}}stCondLst")
    etree.SubElement(inner_stCondLst, f"{{{NSMAP['p']}}}cond",
                     attrib={"delay": str(delay_ms)})
    inner_childTnLst = etree.SubElement(inner_cTn, f"{{{NSMAP['p']}}}childTnLst")

    animEffect = etree.SubElement(inner_childTnLst, f"{{{NSMAP['p']}}}animEffect",
                                  attrib={"transition": "in", "filter": "fade"})
    ae_cBhvr = etree.SubElement(animEffect, f"{{{NSMAP['p']}}}cBhvr")
    etree.SubElement(ae_cBhvr, f"{{{NSMAP['p']}}}cTn", attrib={
        "id": str(next_id + 2), "dur": str(dur_ms)})
    ae_tgtEl = etree.SubElement(ae_cBhvr, f"{{{NSMAP['p']}}}tgtEl")
    etree.SubElement(ae_tgtEl, f"{{{NSMAP['p']}}}spTgt",
                     attrib={"spid": str(shape.shape_id)})


# ── native chart ─────────────────────────────────────────────────────────

def add_chart(slide, chart_type, categories, series_data, l, t, w, h, colors=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)
    chart_frame = slide.shapes.add_chart(chart_type, l, t, w, h, chart_data)
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 70
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10)
    dl.font.color.rgb = WHITE
    dl.font.bold = True

    if colors:
        for series in plot.series:
            if len(series_data) == 1:
                for pt_idx in range(len(categories)):
                    pt = series.points[pt_idx]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = colors[pt_idx % len(colors)]
    chart.chart_style = 2
    return chart_frame


# ── decorative backgrounds ───────────────────────────────────────────────

def deco_hero(slide):
    glow(slide, Inches(11.5), Inches(-0.5), Inches(3.5), TEAL, 92)
    glow(slide, Inches(12.3), Inches(0.8), Inches(2.2), SKY, 94)
    glow(slide, Inches(-0.8), Inches(5.2), Inches(3), CORAL, 93)
    glow(slide, Inches(0.7), Inches(6.5), Inches(1.6), AMBER, 95)
    glow(slide, Inches(6.5), Inches(-1.5), Inches(2.5), VIOLET, 95)
    dot_grid(slide, Inches(10.8), Inches(1.0), 7, 10, Inches(0.24), TEAL, 88)
    dot_grid(slide, Inches(0.5), Inches(5.8), 4, 5, Inches(0.22), CORAL, 90)


def deco_side(slide, side="right"):
    if side == "right":
        glow(slide, Inches(12.5), Inches(0.3), Inches(2.5), TEAL, 93)
        glow(slide, Inches(-0.6), Inches(6.2), Inches(2), VIOLET, 94)
        dot_grid(slide, Inches(11.7), Inches(1.3), 5, 7, Inches(0.22), TEAL, 90)
    else:
        glow(slide, Inches(-1.2), Inches(-0.5), Inches(3), VIOLET, 93)
        glow(slide, Inches(12.8), Inches(5.5), Inches(2.2), AMBER, 94)
        dot_grid(slide, Inches(0.6), Inches(5.8), 4, 5, Inches(0.22), SKY, 90)


# ── SLIDE 1 — TITLE / HERO ───────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_hero(slide)

    # colourful accent bars on top left
    accent_lines(slide, Inches(1.2), Inches(2.1),
                 [TEAL, CORAL, AMBER, VIOLET, SKY])

    # brand tag
    tag = pill(slide, Inches(1.2), Inches(2.3), "PAYSCOPE  •  WEEK 1 PROJECT",
               TEAL, DARK, w=Inches(3.5), h=Inches(0.44))
    anim(slide, tag, 0, 500)

    # massive title
    title = txt(slide, Inches(1.2), Inches(2.85), Inches(11), Inches(1.4),
                "Predict. Explain. Compare.", 54, WHITE, True)
    anim(slide, title, 300, 700)

    # teal highlight subtitle
    sub1 = txt(slide, Inches(1.2), Inches(4.15), Inches(11), Inches(0.6),
               "A data-science salary intelligence platform.", 22, TEAL, False)
    anim(slide, sub1, 600, 500)

    # muted description
    sub2 = txt(slide, Inches(1.2), Inches(4.75), Inches(11), Inches(0.5),
               "ML prediction  •  LLM-generated insight  •  Live analytics dashboard",
               15, MUTED)
    anim(slide, sub2, 900, 500)

    # author info row
    info = txt(slide, Inches(1.2), Inches(6.7), Inches(11), Inches(0.4),
               "Yasser Hamdan     |     AI Engineering Program     |     2026",
               13, MUTED)
    anim(slide, info, 1200, 400)

    # four animated stat pills
    stats = [
        ("607  RECORDS",   TEAL),
        ("DECISION TREE",  AMBER),
        ("FASTAPI + LLM",  CORAL),
        ("SUPABASE LIVE",  VIOLET),
    ]
    for i, (label, color) in enumerate(stats):
        p = pill(slide, Inches(1.2 + i * 2.55), Inches(5.7), label,
                 color, DARK, w=Inches(2.35), h=Inches(0.46))
        anim(slide, p, 1400 + i * 150, 400, "zoom")


# ── SLIDE 2 — PROBLEM → SOLUTION (split) ────────────────────────────────

def slide_problem_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_side(slide, "left")

    # tag + heading
    tag = pill(slide, Inches(0.8), Inches(0.55), "THE GAP WE CLOSE",
               CORAL, DARK, w=Inches(2.8), h=Inches(0.42))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.9),
                "From messy spreadsheets to instant insight.", 30, WHITE, True)
    anim(slide, title, 200, 600)

    # LEFT panel — pain
    left = rect(slide, Inches(0.8), Inches(2.4), Inches(5.8), Inches(4.6),
                SURFACE, 0.06)
    anim(slide, left, 400, 500)

    hard_rect(slide, Inches(0.8), Inches(2.4), Inches(0.12), Inches(4.6), CORAL)
    txt(slide, Inches(1.1), Inches(2.6), Inches(5.4), Inches(0.5),
        "THE PAIN", 13, CORAL, True)
    txt(slide, Inches(1.1), Inches(3.0), Inches(5.4), Inches(0.7),
        "Salary data is broken.", 24, WHITE, True)

    pains = [
        ("Scattered",  "Numbers sit in forums,\nreports & Glassdoor leaks."),
        ("Outdated",   "Most benchmarks are\nmonths or years old."),
        ("No context", "A number alone\ndoesn't explain anything."),
    ]
    for i, (h1, h2) in enumerate(pains):
        y = Inches(4.0 + i * 0.95)
        circle(slide, Inches(1.1), y + Inches(0.15), Inches(0.28), CORAL, 80)
        txt(slide, Inches(1.55), y, Inches(4.9), Inches(0.4), h1, 14, CORAL, True)
        txt(slide, Inches(1.55), y + Inches(0.32), Inches(4.9), Inches(0.65), h2, 11, MUTED)

    # center arrow
    arrow = txt(slide, Inches(6.7), Inches(4.0), Inches(0.8), Inches(1.2),
                "»", 90, TEAL, True, PP_ALIGN.CENTER)
    anim(slide, arrow, 1000, 500, "zoom")

    # RIGHT panel — solution
    right = rect(slide, Inches(7.4), Inches(2.4), Inches(5.1), Inches(4.6),
                 SURFACE, 0.06)
    anim(slide, right, 1200, 500)

    hard_rect(slide, Inches(7.4), Inches(2.4), Inches(0.12), Inches(4.6), TEAL)
    txt(slide, Inches(7.7), Inches(2.6), Inches(4.8), Inches(0.5),
        "PAYSCOPE ANSWERS", 13, TEAL, True)
    txt(slide, Inches(7.7), Inches(3.0), Inches(4.8), Inches(0.7),
        "One click. One answer.", 24, WHITE, True)

    wins = [
        (TEAL,   "Instant ML forecast from a Decision Tree on 607 real records."),
        (SKY,    "LLM written insight that explains *why* the number is what it is."),
        (VIOLET, "Live dashboard of trends: remote, roles, experience, geography."),
    ]
    for i, (col, line) in enumerate(wins):
        y = Inches(4.0 + i * 0.95)
        circle(slide, Inches(7.7), y + Inches(0.18), Inches(0.22), col)
        t1 = txt(slide, Inches(8.1), y, Inches(4.3), Inches(0.85), line, 12, WHITE)
        anim(slide, t1, 1400 + i * 200, 400)


# ── SLIDE 3 — ARCHITECTURE PIPELINE ─────────────────────────────────────

def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_side(slide, "right")

    tag = pill(slide, Inches(0.8), Inches(0.55), "ARCHITECTURE",
               VIOLET, DARK, w=Inches(2.2), h=Inches(0.42))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.8),
                "One request. Five systems. Zero friction.", 28, WHITE, True)
    anim(slide, title, 200, 600)

    sub = txt(slide, Inches(0.8), Inches(1.85), Inches(12), Inches(0.5),
              "Data → model → API → persistence → dashboard. Tracked & explainable.",
              14, MUTED)
    anim(slide, sub, 400, 400)

    # 5 pipeline cards
    steps = [
        ("CSV",        "ds_salaries.csv\n607 records, 11 cols",  TEAL,   "1", "DATA"),
        ("MODEL",      "Decision Tree\nscikit-learn pipeline",    AMBER,  "2", "ML"),
        ("FASTAPI",    "GET /predict\n+ Ollama LLM insight",      CORAL,  "3", "API"),
        ("SUPABASE",   "salary_predictions\nevery call persisted", VIOLET, "4", "DB"),
        ("STREAMLIT",  "Plotly charts\nKPIs + stories",            SKY,    "5", "UI"),
    ]
    box_w = Inches(2.15)
    box_h = Inches(3.0)
    start_x = Inches(0.6)
    y = Inches(2.8)
    gap = Inches(0.27)

    for i, (title_s, body, color, num, kind) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        # shadow
        rect(slide, x + Inches(0.07), y + Inches(0.07), box_w, box_h, SURF_ALT, 0.07)
        # card
        card = rect(slide, x, y, box_w, box_h, SURFACE, 0.07)
        anim(slide, card, 600 + i * 250, 500, "fade")
        # top color strip
        hard_rect(slide, x, y, box_w, Inches(0.16), color)
        # kind tag
        txt(slide, x + Inches(0.25), y + Inches(0.32), box_w - Inches(0.3), Inches(0.3),
            kind, 9, color, True)
        # big number
        txt(slide, x + Inches(0.25), y + Inches(0.55), box_w, Inches(0.9),
            num, 48, color, True)
        # title
        txt(slide, x + Inches(0.25), y + Inches(1.55), box_w - Inches(0.3), Inches(0.45),
            title_s, 17, WHITE, True)
        # divider
        hard_rect(slide, x + Inches(0.25), y + Inches(2.05),
                  Inches(0.4), Inches(0.035), color)
        # body
        txt(slide, x + Inches(0.25), y + Inches(2.15), box_w - Inches(0.3), Inches(0.7),
            body, 10, MUTED)

        # connecting arrow
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.02)
            ar = txt(slide, ax, y + Inches(1.3), Inches(0.25), Inches(0.5),
                     "›", 24, color, True, PP_ALIGN.CENTER)
            anim(slide, ar, 800 + i * 250, 300, "fade")

    # bottom badges
    deploy = [("Local  •  Railway", CORAL), ("Streamlit Cloud", SKY),
              ("Supabase Postgres", VIOLET), ("Ollama llama3.2", AMBER)]
    for i, (label, color) in enumerate(deploy):
        p = pill(slide, Inches(0.8 + i * 3.05), Inches(6.5), label,
                 color, DARK, w=Inches(2.85), h=Inches(0.42))
        anim(slide, p, 2000 + i * 100, 300)


# ── SLIDE 4 — DATA + MODEL ──────────────────────────────────────────────

def slide_data_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_side(slide, "right")

    tag = pill(slide, Inches(0.8), Inches(0.55), "DATA  &  MODEL",
               AMBER, DARK, w=Inches(2.6), h=Inches(0.42))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.8),
                "Real data. Measurable accuracy.", 28, WHITE, True)
    anim(slide, title, 200, 600)

    # LEFT: big stats grid (2x2)
    stats = [
        ("607",      "salary records",   TEAL),
        ("$101K",    "median salary",    AMBER),
        ("11",       "feature columns",  VIOLET),
        ("4 levels", "experience tiers", CORAL),
    ]
    gx = Inches(0.8)
    gy = Inches(2.3)
    gw = Inches(2.85)
    gh = Inches(2.0)
    for i, (big, label, color) in enumerate(stats):
        x = gx + (i % 2) * (gw + Inches(0.25))
        y = gy + (i // 2) * (gh + Inches(0.25))
        card = rect(slide, x, y, gw, gh, SURFACE, 0.08)
        anim(slide, card, 400 + i * 200, 500, "zoom")
        hard_rect(slide, x, y, gw, Inches(0.1), color)
        txt(slide, x + Inches(0.25), y + Inches(0.35), gw - Inches(0.3), Inches(1.0),
            big, 44, color, True)
        txt(slide, x + Inches(0.25), y + Inches(1.35), gw - Inches(0.3), Inches(0.5),
            label.upper(), 11, MUTED, True)

    # RIGHT: native chart of top 5 roles
    chart_bg = rect(slide, Inches(7.1), Inches(2.3), Inches(5.6), Inches(4.25),
                    SURFACE, 0.06)
    anim(slide, chart_bg, 1000, 500)
    txt(slide, Inches(7.35), Inches(2.45), Inches(5.2), Inches(0.4),
        "TOP 5 ROLES BY MEDIAN SALARY (USD)", 10, TEAL, True)

    cats = ["Data\nArchitect", "ML\nEngineer", "Data\nScientist",
            "Data\nEngineer", "Data\nAnalyst"]
    vals = [172_000, 137_000, 128_000, 120_000, 96_000]
    colors = [TEAL, SKY, VIOLET, AMBER, CORAL]
    chart = add_chart(slide, XL_CHART_TYPE.BAR_CLUSTERED,
                      cats, [("Median", vals)],
                      Inches(7.25), Inches(2.85),
                      Inches(5.35), Inches(3.55), colors)
    anim(slide, chart, 1300, 600)

    # bottom model ribbon
    ribbon = rect(slide, Inches(0.8), Inches(6.5), Inches(11.9), Inches(0.65),
                  SURF_ALT, 0.5)
    anim(slide, ribbon, 1500, 400)
    txt(slide, Inches(1.0), Inches(6.6), Inches(11.5), Inches(0.45),
        "MODEL  ·  Decision Tree Regressor    |    Features: experience, role, "
        "remote %, location, company size    |    Trained on 80/20 split",
        12, WHITE, True)


# ── SLIDE 5 — LIVE DEMO (mock window) ───────────────────────────────────

def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_side(slide, "left")

    tag = pill(slide, Inches(0.8), Inches(0.55), "LIVE DEMO",
               SKY, DARK, w=Inches(2.0), h=Inches(0.42))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.8),
                "What the user actually sees.", 28, WHITE, True)
    anim(slide, title, 200, 500)

    # fake browser frame
    frame = rect(slide, Inches(1.0), Inches(2.15), Inches(11.3), Inches(4.9),
                 SURF_ALT, 0.03)
    anim(slide, frame, 400, 600)

    # browser chrome
    hard_rect(slide, Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.4), SURFACE)
    circle(slide, Inches(1.2), Inches(2.25), Inches(0.2), CORAL)
    circle(slide, Inches(1.45), Inches(2.25), Inches(0.2), AMBER)
    circle(slide, Inches(1.7), Inches(2.25), Inches(0.2), LIME)
    urlbar = rect(slide, Inches(2.3), Inches(2.2), Inches(8.5), Inches(0.28),
                  DARK, 0.5)
    txt(slide, Inches(2.55), Inches(2.2), Inches(8.2), Inches(0.3),
        "payscope.streamlit.app",
        10, MUTED, False)

    # inside: hero box
    hero = rect(slide, Inches(1.3), Inches(2.75), Inches(10.7), Inches(1.2),
                SURFACE, 0.04)
    anim(slide, hero, 700, 400)
    circle(slide, Inches(1.55), Inches(3.08), Inches(0.2), TEAL)
    txt(slide, Inches(1.85), Inches(2.9), Inches(5), Inches(0.4),
        "PAYSCOPE", 11, TEAL, True)
    txt(slide, Inches(1.55), Inches(3.25), Inches(10), Inches(0.5),
        "Understand the salary story behind the numbers.", 20, WHITE, True)

    # inside: 4 KPI cards
    kpi_labels = [("$101,570", "MEDIAN", TEAL),
                  ("$142,200", "AVG",    AMBER),
                  ("607",      "RECORDS", VIOLET),
                  ("48%",      "REMOTE",  SKY)]
    for i, (big, lbl, col) in enumerate(kpi_labels):
        x = Inches(1.3 + i * 2.7)
        k = rect(slide, x, Inches(4.1), Inches(2.5), Inches(1.1), SURFACE, 0.08)
        anim(slide, k, 900 + i * 150, 400, "zoom")
        hard_rect(slide, x, Inches(4.1), Inches(2.5), Inches(0.07), col)
        txt(slide, x + Inches(0.2), Inches(4.22), Inches(2.3), Inches(0.5),
            big, 22, col, True)
        txt(slide, x + Inches(0.2), Inches(4.75), Inches(2.3), Inches(0.35),
            lbl, 10, MUTED, True)

    # inside: predict result card
    result = rect(slide, Inches(1.3), Inches(5.4), Inches(5.2), Inches(1.5),
                  SURFACE, 0.06)
    anim(slide, result, 1400, 500, "zoom")
    hard_rect(slide, Inches(1.3), Inches(5.4), Inches(0.08), Inches(1.5), TEAL)
    txt(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.4),
        "PREDICTED SALARY", 10, TEAL, True)
    txt(slide, Inches(1.5), Inches(5.75), Inches(5), Inches(0.7),
        "$119,500", 34, TEAL, True, font="Consolas")
    txt(slide, Inches(1.5), Inches(6.4), Inches(5), Inches(0.4),
        "Senior  ·  ML Engineer  ·  50% remote  ·  US", 11, MUTED)

    # inside: LLM insight card
    insight = rect(slide, Inches(6.7), Inches(5.4), Inches(5.3), Inches(1.5),
                   SURFACE, 0.06)
    anim(slide, insight, 1600, 500, "zoom")
    hard_rect(slide, Inches(6.7), Inches(5.4), Inches(0.08), Inches(1.5), VIOLET)
    txt(slide, Inches(6.9), Inches(5.5), Inches(5), Inches(0.4),
        "LLM INSIGHT", 10, VIOLET, True)
    txt(slide, Inches(6.9), Inches(5.75), Inches(5.1), Inches(1.1),
        "\"This is ~18% above the median. Senior ML\n"
        "talent with hybrid flexibility commands\n"
        "a premium in the US tech market.\"",
        11, WHITE)


# ── SLIDE 6 — KEY DECISIONS ─────────────────────────────────────────────

def slide_decisions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_side(slide, "right")

    tag = pill(slide, Inches(0.8), Inches(0.55), "KEY DECISIONS",
               LIME, DARK, w=Inches(2.4), h=Inches(0.42))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.8),
                "Why I built it this way.", 28, WHITE, True)
    anim(slide, title, 200, 600)

    cards = [
        (TEAL,   "01",  "Decision Tree",
         "Interpretable, fast, handles mixed\ndata types without heavy tuning.\n"
         "Perfect for a Week-1 MVP."),
        (VIOLET, "02",  "Ollama LLM",
         "Runs 100% locally — no API keys,\nno cloud bills, no PII leaks.\n"
         "llama3.2 turns numbers into stories."),
        (SKY,    "03",  "Supabase",
         "Every prediction gets logged.\nFree-tier Postgres + instant REST.\n"
         "Future me can retrain on real usage."),
    ]

    for i, (color, num, title_s, body) in enumerate(cards):
        x = Inches(0.8 + i * 4.15)
        y = Inches(2.4)
        w = Inches(3.95)
        h = Inches(4.1)

        # subtle shadow
        rect(slide, x + Inches(0.08), y + Inches(0.08), w, h, SURF_ALT, 0.08)
        # card
        card = rect(slide, x, y, w, h, SURFACE, 0.08)
        anim(slide, card, 400 + i * 300, 600)

        # top color strip
        hard_rect(slide, x, y, w, Inches(0.14), color)

        # big number
        txt(slide, x + Inches(0.35), y + Inches(0.4), Inches(2), Inches(1.2),
            num, 62, color, True)

        # decorative dots right
        dot_grid(slide, x + w - Inches(1.2), y + Inches(0.5),
                 4, 4, Inches(0.2), color, 88)

        # title
        txt(slide, x + Inches(0.35), y + Inches(1.75), w - Inches(0.4), Inches(0.55),
            title_s, 22, WHITE, True)

        # divider
        hard_rect(slide, x + Inches(0.35), y + Inches(2.35),
                  Inches(0.6), Inches(0.04), color)

        # body
        txt(slide, x + Inches(0.35), y + Inches(2.55), w - Inches(0.4), Inches(1.5),
            body, 13, MUTED)

    # bottom summary ribbon
    ribbon = rect(slide, Inches(0.8), Inches(6.7), Inches(11.75), Inches(0.55),
                  SURF_ALT, 0.5)
    anim(slide, ribbon, 1500, 400)
    txt(slide, Inches(1.0), Inches(6.77), Inches(11.5), Inches(0.4),
        "LESSON LEARNED  ·  Ship simple, ship explainable, ship measurable.",
        12, TEAL, True)


# ── SLIDE 7 — THANK YOU / Q&A ───────────────────────────────────────────

def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    deco_hero(slide)

    # accent lines centered
    accent_lines(slide, Inches(5.4), Inches(2.0),
                 [TEAL, CORAL, AMBER, VIOLET, SKY])

    tag = pill(slide, Inches(5.5), Inches(2.25), "THE END  •  LET'S TALK",
               TEAL, DARK, w=Inches(2.6), h=Inches(0.44))
    anim(slide, tag, 0, 500)

    title = txt(slide, Inches(1.2), Inches(2.85), Inches(11), Inches(1.6),
                "Thank you.", 84, TEAL, True, PP_ALIGN.CENTER)
    anim(slide, title, 300, 800, "zoom")

    sub = txt(slide, Inches(1.2), Inches(4.5), Inches(11), Inches(0.6),
              "Questions, comments, roasts — all welcome.", 22, WHITE, False,
              PP_ALIGN.CENTER)
    anim(slide, sub, 700, 500)

    info = txt(slide, Inches(1.2), Inches(5.15), Inches(11), Inches(0.5),
               "Yasser Hamdan  ·  AI Engineering Program  ·  Week 1",
               14, MUTED, False, PP_ALIGN.CENTER)
    anim(slide, info, 1000, 400)

    # animated bottom pills
    items = [
        ("GITHUB REPO",   TEAL),
        ("LIVE DASHBOARD", VIOLET),
        ("FASTAPI DOCS",  CORAL),
        ("DEMO READY",    AMBER),
    ]
    for i, (label, color) in enumerate(items):
        p = pill(slide, Inches(1.6 + i * 2.6), Inches(6.1), label,
                 color, DARK, w=Inches(2.4), h=Inches(0.5))
        anim(slide, p, 1300 + i * 200, 400, "zoom")


# ── main ─────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)             # 1
    slide_problem_solution(prs)  # 2
    slide_architecture(prs)      # 3
    slide_data_model(prs)        # 4
    slide_demo(prs)              # 5
    slide_decisions(prs)         # 6
    slide_thanks(prs)            # 7

    path = "PayScope_Presentation.pptx"
    prs.save(path)
    print(f"Saved {path} — {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
